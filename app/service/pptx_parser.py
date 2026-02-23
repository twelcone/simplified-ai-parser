"""
PPTX to Markdown Parser

Uses python-pptx to convert PowerPoint presentations to Markdown.
Extracts slide content, tables, and images (as base64).
Speaker notes are excluded (only visible slide content).
"""

import base64
import io
import uuid
import logging
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

logger = logging.getLogger(__name__)

SUPPORTED_IMAGE_FORMATS = ["png", "jpg", "jpeg"]


def _is_supported_image_format(img_format: str) -> bool:
    """Check if image format is supported (png, jpg, jpeg only)."""
    return img_format.lower() in SUPPORTED_IMAGE_FORMATS


def _extract_text_from_shape(shape) -> str:
    """
    Extract text content from a shape.
    Handles text frames with paragraphs.
    """
    if not shape.has_text_frame:
        return ""

    lines = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if text:
            # Check bullet level for indentation
            level = paragraph.level or 0
            indent = "  " * level
            if level > 0:
                lines.append(f"{indent}- {text}")
            else:
                lines.append(text)

    return "\n".join(lines)


def _extract_table_from_shape(shape) -> str:
    """
    Convert a table shape to Markdown table format.
    """
    if not hasattr(shape, 'table'):
        return ""

    table = shape.table
    rows = []

    for row_idx, row in enumerate(table.rows):
        cells = []
        for cell in row.cells:
            # Get cell text, handling merged cells
            cell_text = cell.text.strip() if cell.text else ""
            # Escape pipe characters for markdown
            cell_text = cell_text.replace("|", "\\|").replace("\n", "<br>")
            cells.append(cell_text)
        rows.append(cells)

    if not rows:
        return ""

    # Build markdown table
    lines = []

    # Header row (first row)
    header = rows[0]
    lines.append("| " + " | ".join(header) + " |")

    # Separator row
    lines.append("| " + " | ".join(["---"] * len(header)) + " |")

    # Data rows
    for row in rows[1:]:
        # Ensure row has same number of columns as header
        while len(row) < len(header):
            row.append("")
        lines.append("| " + " | ".join(row[:len(header)]) + " |")

    return "\n".join(lines)


def _extract_image_from_shape(shape) -> dict | None:
    """
    Extract image from a picture shape and convert to base64.
    Returns dict with data_uri and id, or None if unsupported format.
    """
    try:
        image = shape.image
        img_data = image.blob

        # Get image format using PIL
        with io.BytesIO(img_data) as buffer:
            with Image.open(buffer) as pil_img:
                img_format = pil_img.format.lower() if pil_img.format else "png"

        # Skip unsupported formats
        if not _is_supported_image_format(img_format):
            logger.info(f"Skipping unsupported image format: {img_format}")
            return None

        # Convert to base64
        base64_data = base64.b64encode(img_data).decode("utf-8")
        data_uri = f"data:image/{img_format};base64,{base64_data}"

        return {
            "data_uri": data_uri,
            "id": uuid.uuid4().hex[:8]
        }
    except Exception as e:
        logger.warning(f"Failed to extract image: {e}")
        return None


def _get_slide_title(slide) -> str:
    """
    Extract the title from a slide.
    Returns empty string if no title placeholder found.
    """
    try:
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                return title_text
    except Exception:
        pass

    # Fallback: look for a shape that might be a title placeholder
    for shape in slide.shapes:
        try:
            if not shape.has_text_frame:
                continue
            # Only check placeholder_format if shape is a placeholder
            if getattr(shape, 'is_placeholder', False):
                ph_format = shape.placeholder_format
                if ph_format and ph_format.type:
                    type_name = str(ph_format.type)
                    if "TITLE" in type_name:
                        return shape.text.strip()
        except Exception:
            # Skip shapes that cause errors
            continue

    return ""


def _slide_to_markdown(slide, slide_number: int) -> str:
    """
    Convert a single slide to Markdown format.
    Combines title, text content, images, and tables.
    """
    parts = []

    # Get slide title
    title = _get_slide_title(slide)
    if title:
        parts.append(f"## Slide {slide_number}: {title}")
    else:
        parts.append(f"## Slide {slide_number}")

    parts.append("")  # Empty line after heading

    # Track shapes we've already processed (title)
    processed_shapes = set()
    try:
        if slide.shapes.title:
            processed_shapes.add(id(slide.shapes.title))
    except Exception:
        pass

    # Collect content from shapes
    text_content = []
    images = []
    tables = []

    for shape in slide.shapes:
        # Skip already processed shapes (like title)
        if id(shape) in processed_shapes:
            continue

        # Handle pictures (both PICTURE shapes and placeholders with images)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_data = _extract_image_from_shape(shape)
            if img_data:
                images.append(img_data)
            continue

        # Check if placeholder contains an image (image inserted into placeholder)
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            try:
                # Try to access the image property - if it exists, this is a picture placeholder
                if hasattr(shape, 'image') and shape.image:
                    img_data = _extract_image_from_shape(shape)
                    if img_data:
                        images.append(img_data)
                    continue
            except (AttributeError, ValueError):
                # Not a picture placeholder, continue to other checks
                pass

        # Handle tables
        if shape.has_table:
            table_md = _extract_table_from_shape(shape)
            if table_md:
                tables.append(table_md)
            continue

        # Handle grouped shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for sub_shape in shape.shapes:
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img_data = _extract_image_from_shape(sub_shape)
                        if img_data:
                            images.append(img_data)
                    elif hasattr(sub_shape, 'has_text_frame') and sub_shape.has_text_frame:
                        text = _extract_text_from_shape(sub_shape)
                        if text:
                            text_content.append(text)
            except Exception as e:
                logger.warning(f"Failed to process grouped shape: {e}")
            continue

        # Handle text content
        if shape.has_text_frame:
            text = _extract_text_from_shape(shape)
            if text:
                text_content.append(text)

    # Add text content
    if text_content:
        parts.append("\n\n".join(text_content))
        parts.append("")

    # Add images
    for img in images:
        parts.append(f"![image-{img['id']}]({img['data_uri']})")
        parts.append("")

    # Add tables
    for table in tables:
        parts.append(table)
        parts.append("")

    return "\n".join(parts)


def parse_pptx_to_markdown(file_path: str) -> str:
    """
    Convert a PowerPoint file (PPTX) to Markdown format.

    Processing steps:
    1. Load presentation with python-pptx
    2. For each slide:
       - Extract title from title placeholder
       - Extract text from content placeholders
       - Extract tables and convert to Markdown
       - Extract images, filter formats, convert to base64
    3. Combine into formatted Markdown

    Output format per slide:
    ```
    ## Slide N: Title

    Content text...

    ![image-id](data:image/png;base64,...)

    | Col1 | Col2 |
    |------|------|
    | A    | B    |
    ```

    Note: Speaker notes are excluded (only visible slide content).

    Args:
        file_path: Path to the PowerPoint file

    Returns:
        Markdown content as string
    """
    prs = Presentation(file_path)

    markdown_parts = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_md = _slide_to_markdown(slide, slide_number)
        markdown_parts.append(slide_md)

    return "\n".join(markdown_parts).strip()
